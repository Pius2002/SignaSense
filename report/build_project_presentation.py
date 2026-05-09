from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
SOURCE = ROOT / "SignaSense.pptx"
OUTPUT = ROOT / "SignaSense_Project_Presentation.pptx"
MEDIA = ROOT / "pptx_media_extract"
FIGURES = ROOT / "generated_figures"

BACKGROUND = MEDIA / "image1.jpg"
LOGO = MEDIA / "image3.png"

NAVY = RGBColor(20, 31, 54)
BLUE = RGBColor(37, 148, 199)
TEAL = RGBColor(0, 128, 140)
GREEN = RGBColor(19, 130, 70)
MUTED = RGBColor(84, 89, 96)
LIGHT_BLUE = RGBColor(225, 244, 252)
LIGHT_GREEN = RGBColor(226, 244, 234)
LIGHT_YELLOW = RGBColor(255, 244, 206)
WHITE = RGBColor(255, 255, 255)


def delete_slides(prs, zero_based_indices):
    slide_id_list = prs.slides._sldIdLst
    slide_ids = list(slide_id_list)
    for index in sorted(zero_based_indices, reverse=True):
        rel_id = slide_ids[index].rId
        prs.part.drop_rel(rel_id)
        slide_id_list.remove(slide_ids[index])


def reorder_slides(prs, order):
    slide_id_list = prs.slides._sldIdLst
    slide_ids = list(slide_id_list)
    for slide_id in slide_ids:
        slide_id_list.remove(slide_id)
    for index in order:
        slide_id_list.append(slide_ids[index])


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_theme_background(slide, prs, slide_number):
    sw = prs.slide_width
    sh = prs.slide_height
    if BACKGROUND.exists():
        slide.shapes.add_picture(str(BACKGROUND), 0, 0, width=sw, height=sh)
    else:
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, sw, sh)
        rect.fill.solid()
        rect.fill.fore_color.rgb = WHITE
        rect.line.fill.background()

    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(-0.04), Inches(-0.10), height=Inches(1.35))

    add_slide_number(slide, prs, slide_number)


def add_slide_number(slide, prs, slide_number):
    box = slide.shapes.add_textbox(Inches(11.15), Inches(6.92), Inches(1.6), Inches(0.28))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = str(slide_number)
    p.alignment = PP_ALIGN.RIGHT
    run = p.runs[0]
    run.font.name = "Georgia"
    run.font.size = Pt(11)
    run.font.color.rgb = MUTED


def title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.92), Inches(0.37), Inches(11.6), Inches(0.72))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.name = "Georgia"
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = NAVY
    line = slide.shapes.add_connector(1, Inches(3.05), Inches(1.44), Inches(12.5), Inches(1.44))
    line.line.color.rgb = BLUE
    line.line.width = Pt(1.6)
    return box


def subtitle(slide, text, top=1.54):
    box = slide.shapes.add_textbox(Inches(0.94), Inches(top), Inches(10.8), Inches(0.45))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(15)
    run.font.color.rgb = MUTED


def add_bullets(slide, items, x, y, w, h, font_size=19, color=NAVY, gap_pt=6):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(gap_pt)
        p.font.name = "Arial"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
    return box


def add_numbered_steps(slide, steps, x, y, w, item_h=0.72):
    for idx, (label, body) in enumerate(steps, 1):
        top = y + (idx - 1) * (item_h + 0.12)
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(top), Inches(0.44), Inches(0.44))
        circle.fill.solid()
        circle.fill.fore_color.rgb = GREEN
        circle.line.color.rgb = GREEN
        ctf = circle.text_frame
        ctf.clear()
        cp = ctf.paragraphs[0]
        cp.text = str(idx)
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.runs[0]
        cr.font.name = "Arial"
        cr.font.size = Pt(14)
        cr.font.bold = True
        cr.font.color.rgb = WHITE

        head = slide.shapes.add_textbox(Inches(x + 0.58), Inches(top - 0.03), Inches(w - 0.58), Inches(0.26))
        head.text_frame.word_wrap = True
        hp = head.text_frame.paragraphs[0]
        hp.text = label
        hr = hp.runs[0]
        hr.font.name = "Arial"
        hr.font.size = Pt(15)
        hr.font.bold = True
        hr.font.color.rgb = NAVY

        detail = slide.shapes.add_textbox(Inches(x + 0.58), Inches(top + 0.24), Inches(w - 0.58), Inches(0.4))
        detail.text_frame.word_wrap = True
        dp = detail.text_frame.paragraphs[0]
        dp.text = body
        dr = dp.runs[0]
        dr.font.name = "Arial"
        dr.font.size = Pt(12)
        dr.font.color.rgb = MUTED


def add_panel(slide, x, y, w, h, fill_color, line_color=BLUE):
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = fill_color
    panel.line.color.rgb = line_color
    panel.line.width = Pt(1)
    return panel


def add_panel_text(slide, heading, body, x, y, w, h, fill_color):
    add_panel(slide, x, y, w, h, fill_color)
    head = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.16), Inches(w - 0.36), Inches(0.3))
    hp = head.text_frame.paragraphs[0]
    head.text_frame.word_wrap = True
    hp.text = heading
    hr = hp.runs[0]
    hr.font.name = "Arial"
    hr.font.size = Pt(16)
    hr.font.bold = True
    hr.font.color.rgb = NAVY
    body_box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.55), Inches(w - 0.36), Inches(h - 0.62))
    body_box.text_frame.word_wrap = True
    bp = body_box.text_frame.paragraphs[0]
    bp.text = body
    br = bp.runs[0]
    br.font.name = "Arial"
    br.font.size = Pt(12.5)
    br.font.color.rgb = MUTED


def add_picture_fit(slide, image_path, x, y, w, h):
    image_path = Path(image_path)
    if not image_path.exists():
        add_panel_text(slide, "Image needed", image_path.name, x, y, w, h, LIGHT_YELLOW)
        return
    with Image.open(image_path) as img:
        ratio = img.width / img.height
    box_ratio = w / h
    if ratio >= box_ratio:
        pic_w = w
        pic_h = w / ratio
        px = x
        py = y + (h - pic_h) / 2
    else:
        pic_h = h
        pic_w = h * ratio
        px = x + (w - pic_w) / 2
        py = y
    slide.shapes.add_picture(str(image_path), Inches(px), Inches(py), width=Inches(pic_w), height=Inches(pic_h))


def add_image_slide(prs, slide_number, slide_title, lead, image_path, bullets):
    slide = blank_slide(prs)
    add_theme_background(slide, prs, slide_number)
    title(slide, slide_title)
    subtitle(slide, lead)
    add_picture_fit(slide, image_path, 0.75, 1.92, 6.15, 4.45)
    add_bullets(slide, bullets, 7.15, 1.95, 5.45, 4.55, font_size=18)
    return slide


def add_three_panel_slide(prs, slide_number, slide_title, lead, panels):
    slide = blank_slide(prs)
    add_theme_background(slide, prs, slide_number)
    title(slide, slide_title)
    subtitle(slide, lead)
    x_positions = [0.72, 4.62, 8.52]
    fills = [LIGHT_BLUE, LIGHT_GREEN, LIGHT_YELLOW]
    for i, (heading, body) in enumerate(panels):
        add_panel_text(slide, heading, body, x_positions[i], 1.9, 3.25, 4.2, fills[i])
    return slide


def add_functionality_overview_slide(prs, slide_number):
    slide = blank_slide(prs)
    add_theme_background(slide, prs, slide_number)
    title(slide, "System Functionality")
    subtitle(slide, "A simple overview for the live explanation and demonstration.")

    items = [
        ("G", "Smart Glove", "Signs to letters\nand words", GREEN, LIGHT_GREEN),
        ("S", "Smart Stick", "Distance and\nvoice guidance", BLUE, LIGHT_BLUE),
        ("A", "Android App", "Captions, speech\nand control", TEAL, LIGHT_GREEN),
        ("C", "Camera Backup", "Visual signs and\ncustom letters", RGBColor(199, 145, 28), LIGHT_YELLOW),
    ]
    x_positions = [0.85, 3.98, 7.1, 10.2]

    for x, (icon, heading, body, accent, fill_color) in zip(x_positions, items):
        panel = add_panel(slide, x, 2.02, 2.25, 3.75, fill_color, line_color=accent)
        panel.shadow.inherit = False

        circle = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(x + 0.62),
            Inches(2.35),
            Inches(1.0),
            Inches(1.0),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent
        circle.line.color.rgb = accent
        circle_tf = circle.text_frame
        circle_tf.clear()
        circle_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = circle_tf.paragraphs[0]
        p.text = icon
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]
        r.font.name = "Arial"
        r.font.size = Pt(30)
        r.font.bold = True
        r.font.color.rgb = WHITE

        head = slide.shapes.add_textbox(Inches(x + 0.18), Inches(3.68), Inches(1.9), Inches(0.38))
        hp = head.text_frame.paragraphs[0]
        hp.text = heading
        hp.alignment = PP_ALIGN.CENTER
        hr = hp.runs[0]
        hr.font.name = "Arial"
        hr.font.size = Pt(15)
        hr.font.bold = True
        hr.font.color.rgb = NAVY

        detail = slide.shapes.add_textbox(Inches(x + 0.2), Inches(4.26), Inches(1.86), Inches(0.86))
        detail.text_frame.word_wrap = True
        dp = detail.text_frame.paragraphs[0]
        dp.text = body
        dp.alignment = PP_ALIGN.CENTER
        dr = dp.runs[0]
        dr.font.name = "Arial"
        dr.font.size = Pt(13)
        dr.font.color.rgb = MUTED

    note = slide.shapes.add_textbox(Inches(1.2), Inches(6.13), Inches(10.9), Inches(0.44))
    np = note.text_frame.paragraphs[0]
    np.text = "Each part works locally and supports either blind audio guidance or deaf visual communication."
    np.alignment = PP_ALIGN.CENTER
    nr = np.runs[0]
    nr.font.name = "Arial"
    nr.font.size = Pt(15)
    nr.font.color.rgb = NAVY
    return slide


def add_challenges_slide(prs, slide_number):
    slide = blank_slide(prs)
    add_theme_background(slide, prs, slide_number)
    title(slide, "Challenges and Solutions")
    subtitle(slide, "Issues found during development and how the final prototype handled them.")

    rows = [
        ("Flex sensor readings varied", "Added calibration, smoothing, raw-value fallback, and user-defined letter training."),
        ("Power affected glove detection", "Separated testing on USB and direct power; recommended stable 3.3V and proper voltage dividers."),
        ("Bluetooth speaker cannot act as a mic", "Used speaker for output and kept phone voice commands/app input separate."),
        ("Wi-Fi app link caused delay on stick", "Returned stick to local audio guidance with immediate stop and scan commands."),
        ("Camera signs can detect random letters", "Added hold-time filtering, camera backup mode, and train-selected-letter workflow."),
    ]

    left_x = 0.8
    right_x = 6.72
    top = 1.95
    add_panel(slide, left_x, top, 5.55, 0.42, LIGHT_BLUE)
    add_panel(slide, right_x, top, 5.55, 0.42, LIGHT_GREEN)
    for text_value, x in [("Challenge", left_x + 0.18), ("Response", right_x + 0.18)]:
        box = slide.shapes.add_textbox(Inches(x), Inches(top + 0.1), Inches(5.2), Inches(0.25))
        box.text_frame.word_wrap = True
        p = box.text_frame.paragraphs[0]
        p.text = text_value
        r = p.runs[0]
        r.font.name = "Arial"
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = NAVY

    for idx, (challenge, solution) in enumerate(rows):
        y = 2.50 + idx * 0.72
        add_panel(slide, left_x, y, 5.55, 0.53, WHITE, line_color=RGBColor(196, 214, 222))
        add_panel(slide, right_x, y, 5.55, 0.53, WHITE, line_color=RGBColor(196, 214, 222))
        for text_value, x in [(challenge, left_x + 0.18), (solution, right_x + 0.18)]:
            box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.1), Inches(5.2), Inches(0.33))
            box.text_frame.word_wrap = True
            p = box.text_frame.paragraphs[0]
            p.text = text_value
            r = p.runs[0]
            r.font.name = "Arial"
            r.font.size = Pt(12.2)
            r.font.color.rgb = NAVY if x < right_x else MUTED
    return slide


def add_demo_flow_slide(prs, slide_number):
    slide = blank_slide(prs)
    add_theme_background(slide, prs, slide_number)
    title(slide, "Presentation Demonstration Flow")
    subtitle(slide, "A short sequence that shows the system working as one assistive solution.")
    steps = [
        ("Start the app", "Show onboarding, blind/deaf selection, language choice, and green selected-state UI."),
        ("Demonstrate the smart stick", "Bring an obstacle near; the stick says stop, scans left/right, then suggests a safer side."),
        ("Demonstrate the glove", "Connect BLE, make signs, form letters into useful words, and speak committed words."),
        ("Show camera backup", "Use the phone camera when glove power is unavailable or user-defined signs are needed."),
        ("Close with impact", "Explain affordability, offline operation, accessibility, and future field-testing needs."),
    ]
    add_numbered_steps(slide, steps, 0.95, 1.92, 11.2, item_h=0.74)
    return slide


def add_conclusion_slide(prs, slide_number):
    slide = blank_slide(prs)
    add_theme_background(slide, prs, slide_number)
    title(slide, "Conclusion and Recommendations")
    subtitle(slide, "The final prototype demonstrates a practical integrated assistive system.")
    add_panel_text(
        slide,
        "Conclusion",
        "SignaSense combines a smart glove, smart walking stick, and Android app to support both communication and mobility assistance in a low-cost, offline-friendly form.",
        0.82,
        1.95,
        5.5,
        1.8,
        LIGHT_BLUE,
    )
    add_panel_text(
        slide,
        "Recommendations",
        "Improve the glove hardware with stable voltage dividers, expand the sign dataset, package the stick in a durable enclosure, and carry out longer field tests with users.",
        6.7,
        1.95,
        5.5,
        1.8,
        LIGHT_GREEN,
    )
    add_bullets(
        slide,
        [
            "The system supports blind users through audio guidance and deaf users through visual/text output.",
            "The app now supports custom training of any selected letters instead of forcing A-Z training.",
            "Future versions should improve miniaturization, robustness, and language coverage.",
        ],
        1.18,
        4.08,
        10.8,
        1.8,
        font_size=18,
    )
    return slide


def update_slide_numbers(prs):
    for index, slide in enumerate(prs.slides, 1):
        found = False
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            text = shape.text.strip()
            if text.isdigit():
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = str(index)
                p.alignment = PP_ALIGN.RIGHT
                if p.runs:
                    r = p.runs[0]
                    r.font.name = "Georgia"
                    r.font.size = Pt(11)
                    r.font.color.rgb = MUTED
                found = True
        if not found:
            add_slide_number(slide, prs, index)


def add_cover_slide(prs, slide_number):
    slide = blank_slide(prs)
    if BACKGROUND.exists():
        slide.shapes.add_picture(str(BACKGROUND), 0, 0, width=prs.slide_width, height=prs.slide_height)
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(4.1), Inches(0.65), height=Inches(1.05))
    if (MEDIA / "image2.jpeg").exists():
        slide.shapes.add_picture(str(MEDIA / "image2.jpeg"), Inches(9.1), Inches(1.92), height=Inches(2.95))

    main = slide.shapes.add_textbox(Inches(1.55), Inches(1.08), Inches(5.9), Inches(0.82))
    p = main.text_frame.paragraphs[0]
    p.text = "SIGNASENSE"
    r = p.runs[0]
    r.font.name = "Georgia"
    r.font.size = Pt(42)
    r.font.bold = True
    r.font.color.rgb = NAVY

    line = slide.shapes.add_connector(1, Inches(0.95), Inches(2.75), Inches(9.15), Inches(2.75))
    line.line.color.rgb = MUTED
    line.line.width = Pt(1.2)

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(2.98), Inches(8.3), Inches(0.72))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "An Integrated Assistive Solution for the Deaf and Blind"
    sp.alignment = PP_ALIGN.CENTER
    sr = sp.runs[0]
    sr.font.name = "Georgia"
    sr.font.size = Pt(23)
    sr.font.color.rgb = NAVY

    group = slide.shapes.add_textbox(Inches(3.55), Inches(4.05), Inches(3.9), Inches(0.36))
    gp = group.text_frame.paragraphs[0]
    gp.text = "GROUP MEMBERS"
    gp.alignment = PP_ALIGN.CENTER
    gr = gp.runs[0]
    gr.font.name = "Georgia"
    gr.font.size = Pt(18)
    gr.font.bold = True
    gr.font.color.rgb = NAVY

    member_text = (
        "1.   NALUBOWA BENITAH MARGRET      2023/BIT/169/PS\n"
        "2.   ADEKE MARY                                         2023/BIT/034/PS\n"
        "3.   MUHUMUZA MUHAMMAD                 2023/BIT/252/PS\n"
        "4.   MIREMBE JERMIMAH ARNN               2023/BIT/148/PS\n"
        "5.   KANSIIME PIUS                                      2023/BIT/120/PS"
    )
    members = slide.shapes.add_textbox(Inches(0.9), Inches(4.65), Inches(7.9), Inches(1.55))
    mtf = members.text_frame
    mtf.clear()
    mp = mtf.paragraphs[0]
    mp.text = member_text
    mr = mp.runs[0]
    mr.font.name = "Georgia"
    mr.font.size = Pt(18)
    mr.font.color.rgb = NAVY
    add_slide_number(slide, prs, slide_number)
    return slide


def add_text_slide(prs, slide_number, slide_title, bullets, lead=None, image_path=None):
    slide = blank_slide(prs)
    add_theme_background(slide, prs, slide_number)
    title(slide, slide_title)
    if lead:
        subtitle(slide, lead)
    if image_path:
        add_bullets(slide, bullets, 0.92, 1.82, 6.9, 4.95, font_size=18)
        add_picture_fit(slide, image_path, 8.0, 1.65, 4.25, 4.8)
    else:
        add_bullets(slide, bullets, 1.0, 1.86, 11.0, 4.9, font_size=19)
    return slide


def add_objectives_slide(prs, slide_number):
    slide = blank_slide(prs)
    add_theme_background(slide, prs, slide_number)
    title(slide, "Objectives")
    add_panel_text(
        slide,
        "Main Objective",
        "To enhance communication and mobility assistance for the deaf and visually impaired.",
        0.92,
        1.95,
        5.15,
        2.1,
        LIGHT_BLUE,
    )
    add_panel_text(
        slide,
        "Specific Objectives",
        "To examine existing systems that enable communication and mobility assistance for deaf and visually impaired.\n\n"
        "To design and develop a low-cost wearable system that enables communication and mobility assistance for deaf and visually impaired.\n\n"
        "To test and validate the efficiency of the, developed low-cost wearable system that enables communication and mobility assistance for deaf and visually impaired.",
        6.45,
        1.55,
        5.8,
        4.7,
        LIGHT_GREEN,
    )
    return slide


def add_literature_slide(prs, slide_number):
    slide = blank_slide(prs)
    add_theme_background(slide, prs, slide_number)
    title(slide, "Literature Review")
    subtitle(slide, "Existing studies shaped the final system by showing both strengths and remaining gaps.")
    data = [
        ["Authors", "Study", "Strength", "Gap", "How SignaSense bridges it"],
        ["Ambar, Salim et al. (2023)", "Wearable sensor glove for real-time sign language translation", "High ASL alphabet accuracy; real-time display", "Limited vocabulary; no dynamic phrases", "Larger diverse gesture database; dynamic gesture handling"],
        ["Abougarair and Arebi (2022)", "Smart glove for sign language translation", "Simple wearable design; quick response", "No mobility integration; limited phrases", "Expands vocabulary and combines communication with mobility support"],
        ["Olayiwola, Olayode et al. (2025)", "Smart walking stick with voice guidance in an African language", "Voice guidance; culturally relevant for non-English speakers", "Limited obstacle types", "Improves obstacle response through stop, scan, and path guidance"],
        ["Farooq, Shafi et al. (2022)", "IoT enabled intelligent stick for obstacle recognition", "Remote monitoring and navigation support", "Power constraints; no communication integration", "Supports communication features and lower-latency local operation"],
    ]
    table_shape = slide.shapes.add_table(len(data), len(data[0]), Inches(0.55), Inches(1.92), Inches(12.35), Inches(4.65))
    table = table_shape.table
    widths = [1.55, 2.45, 2.45, 2.1, 3.75]
    for idx, width in enumerate(widths):
        table.columns[idx].width = Inches(width)
    for row_idx, row in enumerate(data):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            fill_color = BLUE if row_idx == 0 else (LIGHT_BLUE if row_idx % 2 else WHITE)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_color
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = "Arial"
                paragraph.font.size = Pt(8.2 if row_idx else 8.6)
                paragraph.font.bold = row_idx == 0
                paragraph.font.color.rgb = WHITE if row_idx == 0 else NAVY
    return slide


def add_references_slide(prs, slide_number, title_text, refs):
    slide = blank_slide(prs)
    add_theme_background(slide, prs, slide_number)
    title(slide, title_text)
    add_bullets(slide, refs, 0.78, 1.65, 11.9, 5.1, font_size=9.3, color=NAVY, gap_pt=2)
    return slide


def add_thank_you_slide(prs, slide_number):
    slide = blank_slide(prs)
    if BACKGROUND.exists():
        slide.shapes.add_picture(str(BACKGROUND), 0, 0, width=prs.slide_width, height=prs.slide_height)
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(4.3), Inches(0.55), height=Inches(2.05))
    thanks = slide.shapes.add_textbox(Inches(0.7), Inches(2.55), Inches(11.9), Inches(1.1))
    p = thanks.text_frame.paragraphs[0]
    p.text = "THANK YOU"
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.name = "Georgia"
    r.font.size = Pt(54)
    r.font.bold = True
    r.font.color.rgb = NAVY
    line = slide.shapes.add_connector(1, Inches(3.25), Inches(4.35), Inches(10.9), Inches(4.35))
    line.line.color.rgb = BLUE
    line.line.width = Pt(1.4)
    label = slide.shapes.add_textbox(Inches(4.65), Inches(4.55), Inches(4.0), Inches(0.55))
    lp = label.text_frame.paragraphs[0]
    lp.text = "SIGNASENSE"
    lp.alignment = PP_ALIGN.CENTER
    lr = lp.runs[0]
    lr.font.name = "Georgia"
    lr.font.size = Pt(24)
    lr.font.bold = True
    lr.font.color.rgb = NAVY
    add_slide_number(slide, prs, slide_number)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = 12192000
    prs.slide_height = 6858000

    add_cover_slide(prs, 1)
    add_text_slide(
        prs,
        2,
        "Background",
        [
            "Over 1.5 billion people worldwide live with hearing loss, and about 43 million are blind (WHO, 2024).",
            "In Uganda, around 12% of the population has a disability, yet less than 5% use assistive devices (UBOS, 2023).",
            "Most existing devices address only one disability either hearing or sight.",
            "High costs, limited access, and internet dependency restrict usability in low-resource settings.",
            "There’s an urgent need for affordable, offline, dual-assistive technology to enhance independence and inclusion.",
        ],
        image_path=MEDIA / "image7.png",
    )
    add_text_slide(
        prs,
        3,
        "The Problem Statement",
        [
            "Assistive devices designed for the deaf and the blind are often expensive, require technical expertise and are not designed for our GPS context, making them inaccessible to over 95% of potential users in low-resource settings.",
            "This exclusion limits deaf and visually impaired individuals in Uganda and worldwide from actively participating in activities like education, employment and if better interventions aren’t introduced, this problem will continue to exist.",
            "According to UBOS(2023) and WHO(2024), fewer than 5% of people with visual and hearing impairments have access to assistive tools due to high costs and limited technological suitability.",
            "In line with Uganda’s Vision 2040 and NDP III, this project aims to bridge the assistive technology gap while supporting SDG 3, SDG 8, and SDG 10.",
        ],
        image_path=MEDIA / "image6.jpeg",
    )
    add_objectives_slide(prs, 4)
    add_text_slide(
        prs,
        5,
        "Significance Of The Study",
        [
            "The system empowers deaf and blind individuals to communicate, move, and access services independently, reducing accidents, dependence and social exclusion.",
            "It promotes inclusion in education and workforce, improves quality of life and supports Uganda’s Vision 2040 goal of inclusive human capital development.",
            "The study strengthens our technical and research skills in embedded systems, sensor integration and human computer interaction, preparing us for innovative problem-solving in the tech field.",
            "It provides a valuable foundation for future researchers interested in assistive technologies, guiding them to develop more affordable, adaptive and locally based solutions for people with disabilities.",
        ],
        image_path=MEDIA / "image2.jpeg",
    )
    add_literature_slide(prs, 6)

    next_number = 7
    add_image_slide(
        prs,
        next_number,
            "Integrated System Architecture",
            "The final project connects wearable sensing, obstacle detection, and an Android interface.",
            FIGURES / "figure_02_system_architecture.png",
            [
                "Smart glove reads flex sensors and sends live readings to the Android app through BLE.",
                "Smart stick measures distance with ultrasonic sensing and produces immediate audio guidance.",
                "The Android app provides visual captions, audio interaction, camera backup, and user-defined signs.",
                "Core features are designed to work locally without depending on internet data.",
            ],
    )
    next_number += 1

    add_three_panel_slide(
        prs,
        next_number,
        "Prototype Modules",
        "The project is implemented as three coordinated assistive modules.",
        [
            (
                "Smart Glove",
                "Five flex-sensor inputs capture finger bends. BLE sends raw/bend data to the app, where letters and words are formed.",
            ),
            (
                "Smart Stick",
                "Ultrasonic sensing detects obstacles, speaks distance/guidance, scans left/right, and keeps a buzzer as backup.",
            ),
            (
                "Android App",
                "The app supports blind audio mode, deaf visual mode, smart glove BLE, camera sign backup, and local custom signs.",
            ),
        ],
    )
    next_number += 1

    add_functionality_overview_slide(prs, next_number)
    next_number += 1

    add_image_slide(
        prs,
        next_number,
            "Testing and Validation",
            "The prototype was tested through sensor, app, audio, BLE, and workflow checks.",
            FIGURES / "figure_12_testing_workflow.png",
            [
                "Glove: raw ADC readings, bend response, BLE connection, and letter capture were checked.",
                "Stick: distance accuracy, close-range warnings, scan-left/right logic, and path-clear prompts were tested.",
                "App: phone installation, onboarding, selected-state UI, camera mode, and speech output were verified.",
                "Power behavior was compared between USB power and direct/power-bank operation.",
            ],
    )
    next_number += 1

    add_challenges_slide(prs, next_number)
    next_number += 1

    add_demo_flow_slide(prs, next_number)
    next_number += 1

    add_conclusion_slide(prs, next_number)
    next_number += 1

    refs = [
        "UNDP (2023). A Resilient Future for All: Advancing Disability Inclusion in Disaster Risk Reduction. United Nations Development Programme. Available at: https://www.undp.org/sites/g/files/zskgke326/files/2023-11/undp-a-resilient-future-for-all-advancing-disability-inclusion-in-drr.pdf",
        "WHO (2024). Assistive Technology. World Health Organization. Available at: https://www.who.int/health-topics/assistive-technology",
        "UBOS (2023). Uganda Rapid Assistive Technology Assessment 2023. Uganda Bureau of Statistics. Available at: https://library.health.go.ug/sites/default/files/resources/Uganda%20rapid%20Assistive%20Technology%20Assessment%20Booklet%20V.3-2%20Final.pdf",
        "Ambar, R., Salim, S., Abd Wahab, M.H., Abdul Jamil, M.M. & Ching Phing, T. (2023). Development of a wearable sensor glove for real-time sign language translation. Annals of Emerging Technologies in Computing, 7(5), 25–38. Available at: https://doi.org/10.33166/aetic.2023.05.003",
    ]
    add_references_slide(prs, next_number, "References used (1/2)", refs)
    next_number += 1
    refs_2 = [
        "Abougarair, A.J. & Arebi, W.A. (2022). Smart glove for sign language translation. International Journal of Robotics and Automation, 8(2), 109–117. Available at: https://medcraveonline.com/IRATJ/IRATJ-08-00253.pdf",
        "Olayiwola, A., Olayode, W., Akintayo, T., Oyedeji, A., Olayiwola, D., Osifeko, M. & Abolade, O. (2025). A smart walking stick with voice guidance in an African language for visually impaired persons. Journal of Electrical Systems and Information Technology, 12, 67. Available at: https://doi.org/10.1186/s43067-025-00254-5",
        "Farooq, M.S., Shafi, I., Khan, H., Díez, I.D.L.T., Breñosa, J., Martínez Espinosa, J.C. & Ashraf, I. (2022). IoT enabled intelligent stick for visually impaired people for obstacle recognition. Sensors, 22(22), 8914. Available at: https://doi.org/10.3390/s22228914",
    ]
    add_references_slide(prs, next_number, "References used (2/2)", refs_2)
    next_number += 1
    add_thank_you_slide(prs, next_number)

    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
